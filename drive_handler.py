import datetime
import fitz  # PyMuPDF
from pptx import Presentation
import io


# ------------------ 주차 계산 ------------------ #
def get_current_week(start_date: datetime.date, today: datetime.date) -> int:
    return ((today - start_date).days // 7) + 1


# ------------------ 주차 강의자료 파일 불러오기 ------------------ #
def get_weekly_files(drive, week_num: int, folder_id="YOUR_PARENT_FOLDER_ID"):
    week_name = f"{week_num}주차"

    # 1. 주차 폴더 검색
    folder_list = drive.ListFile({
        'q': f"'{folder_id}' in parents and title contains '{week_name}' and mimeType = 'application/vnd.google-apps.folder' and trashed = false"
    }).GetList()

    if not folder_list:
        return None, None

    week_folder_id = folder_list[0]['id']

    # 2. 폴더 내 강의자료 검색 (PDF, PPTX, DOCX)
    file_list = drive.ListFile({
        'q': f"'{week_folder_id}' in parents and trashed = false"
    }).GetList()

    # 3. 파일 내용 추출 함수 호출
    def find_first_text_file(files):
        for file in files:
            name = file['title'].lower()
            if name.endswith(".pdf") or name.endswith(".pptx") or name.endswith(".docx"):
                return file
        return None

    # 이번 주차 + 지난 주차 기준
    this_week_file = find_first_text_file(file_list)

    # 지난 주차 자료도 함께 불러오기
    last_week_file = None
    if week_num > 1:
        last_folder_name = f"{week_num - 1}주차"
        last_folders = drive.ListFile({
            'q': f"'{folder_id}' in parents and title contains '{last_folder_name}' and mimeType = 'application/vnd.google-apps.folder' and trashed = false"
        }).GetList()
        if last_folders:
            last_folder_id = last_folders[0]['id']
            last_files = drive.ListFile({
                'q': f"'{last_folder_id}' in parents and trashed = false"
            }).GetList()
            last_week_file = find_first_text_file(last_files)

    # 텍스트 추출
    last_text = extract_text_from_file(last_week_file) if last_week_file else None
    this_text = extract_text_from_file(this_week_file) if this_week_file else None

    return last_text, this_text


# ------------------ 파일 내용 텍스트 추출 ------------------ #
def extract_text_from_file(file):
    if not file:
        return None

    file_name = file['title'].lower()
    file.GetContentFile(file_name)

    if file_name.endswith(".pdf"):
        doc = fitz.open(file_name)
        text = "".join([page.get_text() for page in doc])
        return text

    elif file_name.endswith(".pptx"):
        prs = Presentation(file_name)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    elif file_name.endswith(".docx"):
        import docx
        doc = docx.Document(file_name)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text

    return None
